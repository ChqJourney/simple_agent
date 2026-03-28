from __future__ import annotations

import re
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover - exercised via runtime validation
    Presentation = None


MAX_PPTX_FILE_BYTES = 20 * 1024 * 1024
MAX_TITLE_CHARS = 80


def _require_python_pptx() -> None:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed. Install `python-pptx` to use PowerPoint document tools.")


def _normalize_text(text: Any) -> str:
    return " ".join(str(text or "").split()).strip()


def _make_anchor(text: str) -> str:
    anchor = re.sub(r"[^\w\s-]", "", text.lower()).strip()
    anchor = re.sub(r"[\s_]+", "-", anchor)
    return anchor


def _truncate_title(text: str) -> str:
    normalized = _normalize_text(text)
    if len(normalized) <= MAX_TITLE_CHARS:
        return normalized
    return normalized[:MAX_TITLE_CHARS].rstrip() + "..."


class PptxReader:
    def __init__(self, pptx_path: str | Path):
        _require_python_pptx()
        self.pptx_path = Path(pptx_path)
        if not self.pptx_path.exists():
            raise FileNotFoundError(f"PowerPoint file not found: {self.pptx_path}")
        file_size = self.pptx_path.stat().st_size
        if file_size > MAX_PPTX_FILE_BYTES:
            raise ValueError(f"File too large: {file_size} bytes (max: {MAX_PPTX_FILE_BYTES} bytes)")

        self.presentation = Presentation(self.pptx_path)
        self.slides: list[dict[str, Any]] = [self._serialize_slide(slide, index + 1) for index, slide in enumerate(self.presentation.slides)]

    @property
    def slide_count(self) -> int:
        return len(self.slides)

    def _serialize_slide(self, slide: Any, slide_number: int) -> dict[str, Any]:
        title = ""
        if getattr(slide.shapes, "title", None) is not None:
            title = _normalize_text(getattr(slide.shapes.title, "text", ""))

        text_blocks: list[dict[str, Any]] = []
        for shape_index, shape in enumerate(getattr(slide, "shapes", []), start=1):
            if not getattr(shape, "has_text_frame", False):
                continue
            text = _normalize_text(getattr(shape, "text", ""))
            if not text:
                continue
            text_blocks.append(
                {
                    "shape_index": shape_index,
                    "text": text,
                }
            )
            if not title:
                title = text

        notes_text = ""
        notes_slide = getattr(slide, "notes_slide", None)
        if notes_slide is not None:
            note_parts = []
            for shape in getattr(notes_slide, "shapes", []):
                if not getattr(shape, "has_text_frame", False):
                    continue
                text = _normalize_text(getattr(shape, "text", ""))
                if text:
                    note_parts.append(text)
            notes_text = "\n".join(note_parts)

        if not title:
            title = f"Slide {slide_number}"

        return {
            "slide_number": slide_number,
            "title": _truncate_title(title),
            "text_blocks": text_blocks,
            "notes_text": notes_text,
            "text": "\n".join(block["text"] for block in text_blocks),
        }

    def get_structure(self, max_nodes: int = 200) -> dict[str, Any]:
        items = [
            {
                "title": slide["title"],
                "level": 1,
                "anchor": _make_anchor(f"slide-{slide['slide_number']}-{slide['title']}"),
                "locator": {
                    "slide_number": int(slide["slide_number"]),
                },
                "slide_summary": {
                    "text_block_count": len(slide["text_blocks"]),
                    "has_notes": bool(slide["notes_text"]),
                },
            }
            for slide in self.slides[:max_nodes]
        ]
        return {
            "pptx_path": str(self.pptx_path),
            "slide_count": self.slide_count,
            "structure_type": "pptx_slide_map",
            "items": items,
        }

    def search(
        self,
        query: str,
        *,
        mode: str = "plain",
        case_sensitive: bool = False,
        max_results: int = 50,
    ) -> dict[str, Any]:
        flags = 0 if case_sensitive else re.IGNORECASE
        pattern = re.compile(query if mode == "regex" else re.escape(query), flags)
        items: list[dict[str, Any]] = []

        for slide in self.slides:
            for block in slide["text_blocks"]:
                if len(items) >= max_results:
                    break
                text = str(block["text"])
                match = pattern.search(text)
                if not match:
                    continue
                items.append(
                    {
                        "source_type": "slide_text",
                        "slide_number": int(slide["slide_number"]),
                        "slide_title": str(slide["title"]),
                        "shape_index": int(block["shape_index"]),
                        "match_text": match.group(0),
                        "text": text,
                        "context_before": "",
                        "context_after": "",
                    }
                )
            if len(items) >= max_results:
                break

            notes_text = str(slide["notes_text"] or "")
            if notes_text and len(items) < max_results:
                match = pattern.search(notes_text)
                if match:
                    items.append(
                        {
                            "source_type": "slide_notes",
                            "slide_number": int(slide["slide_number"]),
                            "slide_title": str(slide["title"]),
                            "match_text": match.group(0),
                            "text": notes_text,
                            "context_before": "",
                            "context_after": "",
                        }
                    )
            if len(items) >= max_results:
                break

        return {
            "pptx_path": str(self.pptx_path),
            "slide_count": self.slide_count,
            "items": items,
        }

    def read_slides(
        self,
        slide_start: int,
        slide_end: int,
    ) -> dict[str, Any]:
        if slide_start < 1 or slide_end < slide_start:
            raise ValueError("Invalid slide range: slide_end must be >= slide_start and slide_start must be >= 1")
        if slide_end > self.slide_count:
            raise ValueError(f"Slide range {slide_start}-{slide_end} exceeds slide count {self.slide_count}")

        items = []
        for slide in self.slides[slide_start - 1:slide_end]:
            items.append(
                {
                    "slide_number": int(slide["slide_number"]),
                    "title": str(slide["title"]),
                    "text": str(slide["text"]),
                    "notes_text": str(slide["notes_text"]),
                    "text_blocks": list(slide["text_blocks"]),
                }
            )
        return {
            "pptx_path": str(self.pptx_path),
            "slide_start": slide_start,
            "slide_end": slide_end,
            "items": items,
        }


def get_pptx_structure(pptx_path: str | Path, max_nodes: int = 200) -> dict[str, Any]:
    reader = PptxReader(pptx_path)
    return reader.get_structure(max_nodes=max_nodes)


def search_pptx_document(
    pptx_path: str | Path,
    query: str,
    *,
    mode: str = "plain",
    case_sensitive: bool = False,
    max_results: int = 50,
) -> dict[str, Any]:
    reader = PptxReader(pptx_path)
    return reader.search(
        query,
        mode=mode,
        case_sensitive=case_sensitive,
        max_results=max_results,
    )


def read_pptx_slides(
    pptx_path: str | Path,
    slide_start: int,
    slide_end: int,
) -> dict[str, Any]:
    reader = PptxReader(pptx_path)
    return reader.read_slides(slide_start, slide_end)
